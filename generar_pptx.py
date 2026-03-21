
import io
import json
import math
import os
import sys
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ── Unidades: todo en pulgadas ──────────────
W  = 10.0   # ancho slide
H  = 5.625  # alto slide

def rgb(hex6: str):
    """'065A82' → RGBColor"""
    h = hex6.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def add_rect(slide, x, y, w, h, fill_hex, alpha=255):
    """Añade un rectángulo de color sólido."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.fill.background()   # sin borde
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    return shape

def add_oval(slide, x, y, w, h, fill_hex):
    """Añade un óvalo."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        9,  # oval
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=14, bold=False, italic=False,
             color_hex="000000", font_name="Calibri",
             align=PP_ALIGN.LEFT, valign="middle", wrap=True):
    """Añade un cuadro de texto."""
    from pptx.enum.text import MSO_ANCHOR
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap

    # alineación vertical
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.TOP

    # margen interno 0
    tf.margin_left   = Emu(0)
    tf.margin_right  = Emu(0)
    tf.margin_top    = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size      = Pt(font_size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = rgb(color_hex)
    run.font.name      = font_name
    return txBox

def add_notes(slide, text):
    """Añade notas del presentador."""
    if not text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text

# ── Iconos: círculos de colores con letra inicial ─────────────────────
# Sin dependencias externas: dibujamos un círculo PNG con Pillow

# Letras simples ASCII - funcionan con cualquier fuente en cualquier SO
ICON_LETTERS = {
    "FaBook": "B",   "FaLightbulb": "i",  "FaChartBar": "G",
    "FaStar": "*",   "FaUsers": "U",       "FaCog": "C",
    "FaFlask": "F",  "FaGlobe": "O",       "FaHeart": "H",
    "FaSearch": "S", "FaCheckCircle": "V", "FaBrain": "P",
    "FaCalculator": "N", "FaAtom": "A",    "FaDna": "D",
    "FaLeaf": "L",   "FaCode": "{ }",      "FaHistory": "T",
    "FaBalanceScale": "J", "FaMusic": "M", "FaLock": "K",
    "FaArrowRight": "R", "FaSun": "Q",     "FaThermometerHalf": "E",
    "FaTint": "W",   "FaWind": "Z",        "FaCircle": "o",
}

# Rutas de fuentes en orden de preferencia (Windows primero, luego Linux/Mac)
_FONT_CANDIDATES = [
    "C:/Windows/Fonts/arialbd.ttf",
    "C:/Windows/Fonts/calibrib.ttf",
    "C:/Windows/Fonts/verdanab.ttf",
    "/System/Library/Fonts/Helvetica.ttc",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    "/usr/share/fonts/truetype/crosextra/Caladea-Bold.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
]

def _get_font(size_pt: int):
    """Devuelve la mejor fuente disponible en el sistema."""
    from PIL import ImageFont
    for path in _FONT_CANDIDATES:
        try:
            return ImageFont.truetype(path, size_pt)
        except Exception:
            pass
    return ImageFont.load_default()

def make_icon_png(icon_name: str, color_hex: str, size: int = 256) -> bytes:
    """Genera un PNG circular con letra usando Pillow. Funciona en cualquier SO."""
    letra = ICON_LETTERS.get(icon_name, (icon_name[2] if len(icon_name) > 2 else "?"))
    img   = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    draw  = ImageDraw.Draw(img)
    h     = color_hex.lstrip("#")
    fill  = (int(h[0:2],16), int(h[2:4],16), int(h[4:6],16), 220)
    draw.ellipse([6, 6, size-6, size-6], fill=fill)
    font  = _get_font(int(size * 0.42))
    bbox  = draw.textbbox((0, 0), letra, font=font)
    tw    = bbox[2] - bbox[0]
    th    = bbox[3] - bbox[1]
    draw.text(((size - tw) // 2 - bbox[0], (size - th) // 2 - bbox[1]),
              letra, fill=(255, 255, 255, 245), font=font)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()

def add_icon(slide, icon_name, color_hex, x, y, size_in):
    """Añade un icono como imagen PNG generada con Pillow."""
    png_bytes = make_icon_png(icon_name, color_hex, 256)
    buf = io.BytesIO(png_bytes)
    slide.shapes.add_picture(buf, Inches(x), Inches(y), Inches(size_in), Inches(size_in))

# ── GENERADOR PRINCIPAL ───────────────────────────────────────────────

def generar_presentacion(estructura: dict, paleta: dict, ruta_salida: str):
    prs = Presentation()
    prs.slide_width  = Inches(W)
    prs.slide_height = Inches(H)

    # Colores de la paleta
    C = {k: v.lstrip("#") for k, v in paleta.items() if k != "fondo_claro"}
    C["white"] = "FFFFFF"
    C["gray"]  = "64748B"
    es_claro   = paleta.get("fondo_claro", False)
    C["light"] = "F8F9FA" if not es_claro else paleta.get("dark","F8F9FA").lstrip("#")

    FTITULO = "Georgia"
    FBODY   = "Calibri"

    blank_layout = prs.slide_layouts[6]  # layout en blanco

    diapositivas = estructura.get("diapositivas", [])

    for i, diapo in enumerate(diapositivas):
        slide = prs.slides.add_slide(blank_layout)
        tipo  = diapo.get("tipo", "contenido")

        # ── PORTADA ───────────────────────────────────────────────
        if tipo == "portada":
            add_rect(slide, 0, 0, W, H, C["dark"])
            if es_claro:
                # Fondo claro: barra superior con color primario
                add_rect(slide, 0, 0, W, 2.2, C["primary"])
                add_rect(slide, 0, 0, 0.22, H, C["primary"])
                add_rect(slide, 0.22, 0, 0.08, H, C["secondary"])
                titulo_color  = C["white"]
                sub_color     = C["dark"]   # texto sobre fondo claro
                asig_color    = C["gray"]
                # Línea separadora
                add_rect(slide, 0.6, 2.1, 8.8, 0.06, C["secondary"])
            else:
                # Fondo oscuro: franjas laterales clásicas
                add_rect(slide, 0, 0, 0.22, H, C["primary"])
                add_rect(slide, 0.22, 0, 0.08, H, C["secondary"])
                add_rect(slide, 0.6, 1.5, 8.8, 0.06, C["primary"])
                titulo_color  = C["white"]
                sub_color     = C["secondary"]
                asig_color    = C["gray"]

            titulo_pres = estructura.get("titulo_presentacion") or diapo.get("titulo", "")
            titulo_y    = 0.3 if es_claro else 0.5
            add_text(slide, titulo_pres,
                     0.6, titulo_y, 8.8, 1.3,
                     font_size=38, bold=True, color_hex=titulo_color,
                     font_name=FTITULO, valign="bottom")

            subtitulo = diapo.get("subtitulo", "")
            if subtitulo:
                sub_y = 1.7 if es_claro else 1.65
                add_text(slide, subtitulo,
                         0.6, sub_y, 8.8, 1.1,
                         font_size=20, italic=True, color_hex=sub_color,
                         font_name=FBODY, valign="top")

            asignatura = estructura.get("asignatura", "")
            if asignatura:
                add_text(slide, asignatura,
                         0.6, 4.9, 8.8, 0.5,
                         font_size=13, color_hex=asig_color,
                         font_name=FBODY, valign="middle")

            add_notes(slide, diapo.get("notas", ""))
            continue

        # ── ENCABEZADO COMÚN (todas las slides excepto portada) ───
        # Fondo claro
        add_rect(slide, 0, 0, W, H, C["light"])
        # Barra superior
        add_rect(slide, 0, 0, W, 1.15, C["primary"])
        # Franja lateral
        add_rect(slide, 0, 0, 0.15, H, C["primary"])
        # Número de slide
        add_text(slide, str(i),
                 9.1, 0.25, 0.7, 0.5,
                 font_size=11, color_hex=C["secondary"],
                 font_name=FBODY, align=PP_ALIGN.RIGHT)

        # ── RESUMEN ───────────────────────────────────────────────
        if tipo == "resumen":
            add_text(slide, diapo.get("titulo", "Puntos clave"),
                     0.35, 0.2, 9.3, 0.7,
                     font_size=28, bold=True, color_hex=C["white"],
                     font_name=FTITULO, valign="middle")

            for idx, punto in enumerate(diapo.get("puntos", [])):
                y = 1.4 + idx * 0.72
                add_rect(slide, 0.35, y + 0.05, 0.06, 0.35, C["primary"])
                add_text(slide, punto,
                         0.55, y, 9.1, 0.5,
                         font_size=15, color_hex=C["texto"],
                         font_name=FBODY, valign="middle")

            add_notes(slide, diapo.get("notas", ""))
            continue

        # ── CONTENIDO ─────────────────────────────────────────────
        if tipo == "contenido":
            add_text(slide, diapo.get("titulo", ""),
                     0.35, 0.2, 8.6, 0.7,
                     font_size=26, bold=True, color_hex=C["white"],
                     font_name=FTITULO, valign="middle")

            icono = diapo.get("icono", "FaCircle")
            add_oval(slide, 7.6, 1.25, 2.0, 2.0, C["secondary"])
            add_icon(slide, icono, C["primary"], 7.85, 1.5, 1.5)

            for idx, punto in enumerate(diapo.get("puntos", [])):
                y = 1.35 + idx * 0.8
                add_oval(slide, 0.3, y + 0.04, 0.38, 0.38, C["primary"])
                add_text(slide, str(idx + 1),
                         0.3, y + 0.04, 0.38, 0.38,
                         font_size=12, bold=True, color_hex=C["white"],
                         font_name=FBODY, align=PP_ALIGN.CENTER, valign="middle")
                add_text(slide, punto,
                         0.82, y, 6.6, 0.5,
                         font_size=14, color_hex=C["texto"],
                         font_name=FBODY, valign="middle")

            add_notes(slide, diapo.get("notas", ""))
            continue

        # ── DOS COLUMNAS ──────────────────────────────────────────
        if tipo == "dos_columnas":
            add_text(slide, diapo.get("titulo", ""),
                     0.35, 0.2, 9.3, 0.7,
                     font_size=26, bold=True, color_hex=C["white"],
                     font_name=FTITULO, valign="middle")

            cols  = [diapo.get("columna_izq"), diapo.get("columna_der")]
            col_x = [0.3, 5.25]

            for ci, col in enumerate(cols):
                if not col:
                    continue
                x = col_x[ci]
                # Tarjeta blanca
                add_rect(slide, x, 1.25, 4.65, 4.1, C["white"])
                # Encabezado de columna
                hdr_color = C["primary"] if ci == 0 else C["secondary"]
                add_rect(slide, x, 1.25, 4.65, 0.42, hdr_color)
                txt_color = C["white"] if ci == 0 else C["dark"]
                add_text(slide, col.get("titulo", ""),
                         x + 0.15, 1.25, 4.35, 0.42,
                         font_size=13, bold=True, color_hex=txt_color,
                         font_name=FBODY, valign="middle")

                for ii, item in enumerate(col.get("items", [])):
                    y = 1.82 + ii * 0.65
                    add_rect(slide, x + 0.18, y + 0.1, 0.06, 0.3, C["primary"])
                    add_text(slide, item,
                             x + 0.36, y, 4.0, 0.5,
                             font_size=13, color_hex=C["texto"],
                             font_name=FBODY, valign="middle")

            add_notes(slide, diapo.get("notas", ""))
            continue

        # ── ESTADÍSTICA ───────────────────────────────────────────
        if tipo == "estadistica":
            add_text(slide, diapo.get("titulo", ""),
                     0.35, 0.2, 9.3, 0.7,
                     font_size=26, bold=True, color_hex=C["white"],
                     font_name=FTITULO, valign="middle")

            datos  = diapo.get("datos", [])
            n      = len(datos)
            tarj_w = 4.3 if n == 2 else 2.85
            start_x= 0.75 if n == 2 else 0.3
            gap    = 4.6  if n == 2 else 3.2

            for di, dato in enumerate(datos):
                x = start_x + di * gap
                card_color = C["primary"] if di % 2 == 0 else C["secondary"]
                add_rect(slide, x, 1.4, tarj_w, 2.9, card_color)
                val_color = C["white"]
                add_text(slide, dato.get("valor", ""),
                         x, 1.7, tarj_w, 1.4,
                         font_size=48, bold=True, color_hex=val_color,
                         font_name=FTITULO, align=PP_ALIGN.CENTER, valign="middle")
                add_text(slide, dato.get("etiqueta", ""),
                         x + 0.15, 3.1, tarj_w - 0.3, 0.9,
                         font_size=14, color_hex=C["white"],
                         font_name=FBODY, align=PP_ALIGN.CENTER, valign="top")

            add_notes(slide, diapo.get("notas", ""))
            continue

        # ── LISTA DE ICONOS ───────────────────────────────────────
        if tipo == "lista_iconos":
            add_text(slide, diapo.get("titulo", ""),
                     0.35, 0.2, 9.3, 0.7,
                     font_size=26, bold=True, color_hex=C["white"],
                     font_name=FTITULO, valign="middle")

            items   = diapo.get("items", [])[:4]
            n_items = len(items)
            cols    = 2 if n_items > 2 else n_items
            rows    = math.ceil(n_items / cols) if cols else 1
            card_w  = 4.55 if cols == 2 else 9.3
            card_h  = 1.85 if rows == 2 else 3.8

            for ii, item in enumerate(items):
                col = ii % cols
                row = ii // cols
                x   = 0.3 + col * (card_w + 0.2)
                y   = 1.4 + row * (card_h + 0.15)

                add_rect(slide, x, y, card_w, card_h, C["white"])
                stripe_color = C["primary"] if ii % 2 == 0 else C["secondary"]
                add_rect(slide, x, y, 0.08, card_h, stripe_color)

                icon_color = C["primary"] if ii % 2 == 0 else C["secondary"]
                icon_size  = 0.65 if card_h > 2 else 0.5
                add_icon(slide, item.get("icono", "FaCircle"),
                         icon_color, x + 0.2, y + (card_h - icon_size) / 2, icon_size)

                add_text(slide, item.get("titulo", ""),
                         x + 1.05, y + 0.18, card_w - 1.2, 0.42,
                         font_size=14, bold=True, color_hex=C["texto"],
                         font_name=FBODY, valign="middle")
                add_text(slide, item.get("descripcion", ""),
                         x + 1.05, y + 0.6, card_w - 1.2, card_h - 0.75,
                         font_size=12, color_hex=C["gray"],
                         font_name=FBODY, valign="top")

            add_notes(slide, diapo.get("notas", ""))
            continue

        # ── TIPO DESCONOCIDO: título simple ───────────────────────
        add_text(slide, diapo.get("titulo", "(sin título)"),
                 0.35, 0.2, 9.3, 0.7,
                 font_size=26, bold=True, color_hex=C["white"],
                 font_name=FTITULO, valign="middle")
        add_notes(slide, diapo.get("notas", ""))

    prs.save(ruta_salida)
    print(f"✅  PPTX guardado -> {ruta_salida}")


# ── ENTRY POINT ───────────────────────────────────────────────────────
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Uso: python generar_pptx.py datos.json")
        sys.exit(1)

    with open(sys.argv[1], encoding="utf-8") as f:
        data = json.load(f)

    generar_presentacion(data["estructura"], data["paleta"], data["salida"])
